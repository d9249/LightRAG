"""
파일 처리 유틸리티
"""
import os
import uuid
import magic
import mimetypes
from pathlib import Path
from typing import List, Optional, Tuple
from app.config import config


class FileHandler:
    """파일 처리 유틸리티 클래스"""
    
    def __init__(self):
        self.supported_extensions = config.supported_extensions_list
        self.max_file_size = config.max_file_size
        self.inputs_dir = Path(config.inputs_dir)
        self.inputs_dir.mkdir(parents=True, exist_ok=True)
    
    def is_supported_file(self, filename: str) -> bool:
        """지원되는 파일 형식인지 확인"""
        file_ext = Path(filename).suffix.lower()
        return file_ext in self.supported_extensions
    
    def is_valid_file_size(self, file_size: int) -> bool:
        """파일 크기가 유효한지 확인"""
        return file_size <= self.max_file_size
    
    def generate_unique_filename(self, original_filename: str) -> str:
        """고유한 파일명 생성"""
        file_ext = Path(original_filename).suffix
        unique_id = str(uuid.uuid4())
        return f"{unique_id}{file_ext}"
    
    def get_mime_type(self, file_path: str) -> str:
        """파일의 MIME 타입 추출"""
        try:
            # python-magic을 사용하여 파일 타입 감지
            mime_type = magic.from_file(file_path, mime=True)
            return mime_type
        except Exception:
            # fallback: mimetypes 모듈 사용
            mime_type, _ = mimetypes.guess_type(file_path)
            return mime_type or "application/octet-stream"
    
    def save_uploaded_file(self, file_content: bytes, original_filename: str) -> Tuple[str, str]:
        """업로드된 파일 저장
        
        Returns:
            Tuple[저장된_파일_경로, 고유_파일명]
        """
        unique_filename = self.generate_unique_filename(original_filename)
        file_path = self.inputs_dir / unique_filename
        
        with open(file_path, "wb") as f:
            f.write(file_content)
        
        return str(file_path), unique_filename
    
    def delete_file(self, file_path: str) -> bool:
        """파일 삭제"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                return True
            return False
        except Exception:
            return False
    
    def get_file_size(self, file_path: str) -> int:
        """파일 크기 반환"""
        try:
            return os.path.getsize(file_path)
        except Exception:
            return 0
    
    def scan_directory(self, directory_path: str, recursive: bool = True, 
                      extensions: Optional[List[str]] = None) -> List[str]:
        """디렉터리 스캔하여 지원되는 파일 목록 반환"""
        if extensions is None:
            extensions = self.supported_extensions
        
        found_files = []
        directory = Path(directory_path)
        
        if not directory.exists():
            return found_files
        
        pattern = "**/*" if recursive else "*"
        
        for file_path in directory.glob(pattern):
            if file_path.is_file():
                file_ext = file_path.suffix.lower()
                if file_ext in extensions:
                    found_files.append(str(file_path))
        
        return found_files
    
    def extract_text_content(self, file_path: str) -> str:
        """파일에서 텍스트 내용 추출"""
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.txt':
                return self._extract_from_txt(file_path)
            elif file_ext == '.md':
                return self._extract_from_markdown(file_path)
            elif file_ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._extract_from_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._extract_from_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_from_pptx(file_path)
            elif file_ext == '.json':
                return self._extract_from_json(file_path)
            elif file_ext == '.xml':
                return self._extract_from_xml(file_path)
            elif file_ext == '.csv':
                return self._extract_from_csv(file_path)
            else:
                # 기본적으로 텍스트 파일로 시도
                return self._extract_from_txt(file_path)
        except Exception as e:
            raise Exception(f"텍스트 추출 실패: {str(e)}")
    
    def _extract_from_txt(self, file_path: str) -> str:
        """텍스트 파일에서 내용 추출"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_from_markdown(self, file_path: str) -> str:
        """마크다운 파일에서 내용 추출"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """PDF 파일에서 텍스트 추출"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except ImportError:
            raise Exception("PDF 처리를 위해 PyPDF2가 필요합니다")
    
    def _extract_from_docx(self, file_path: str) -> str:
        """DOCX 파일에서 텍스트 추출"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            raise Exception("DOCX 처리를 위해 python-docx가 필요합니다")
    
    def _extract_from_excel(self, file_path: str) -> str:
        """Excel 파일에서 텍스트 추출"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            return text
        except ImportError:
            raise Exception("Excel 처리를 위해 openpyxl이 필요합니다")
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            raise Exception("PowerPoint 처리를 위해 python-pptx가 필요합니다")
    
    def _extract_from_json(self, file_path: str) -> str:
        """JSON 파일에서 텍스트 추출"""
        import json
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
    
    def _extract_from_xml(self, file_path: str) -> str:
        """XML 파일에서 텍스트 추출"""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            def extract_text_from_element(element):
                text = element.text or ""
                for child in element:
                    text += extract_text_from_element(child)
                text += element.tail or ""
                return text
            
            return extract_text_from_element(root)
        except Exception:
            # XML 파싱 실패 시 원본 텍스트 반환
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
    
    def _extract_from_csv(self, file_path: str) -> str:
        """CSV 파일에서 텍스트 추출"""
        import csv
        text = ""
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                text += "\t".join(row) + "\n"
        return text